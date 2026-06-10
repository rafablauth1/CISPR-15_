# -*- coding: utf-8 -*-
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

SHOTS = os.path.join(os.environ.get("TEMP", "/tmp"), "cispr15-shots")
ROOT  = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ICON  = os.path.join(ROOT, "assets", "icon.png")
PUCRS = os.path.join(ROOT, "public", "formularios", "emc", "pucrs-logo.png")
OUTDIR = os.path.join(ROOT, "releases")
OUT   = os.path.join(OUTDIR, "CISPR15-LABELO - Visao Geral do App.pptx")

# Paleta (tema do app)
BG    = RGBColor(0x0B, 0x0F, 0x17)
PANEL = RGBColor(0x12, 0x18, 0x24)
GOLD  = RGBColor(0xE6, 0xB7, 0x3C)
WHITE = RGBColor(0xF2, 0xF4, 0xF8)
GRAY  = RGBColor(0xA6, 0xB0, 0xC0)
LINE  = RGBColor(0x26, 0x30, 0x40)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp

def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=6):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    first = True
    for item in runs:
        s, size, color, bold = (list(item) + [False])[:4]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.space_after = Pt(sp_after); p.space_before = Pt(0)
        r = p.add_run(); r.text = s
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
        r.font.name = "Segoe UI"
    return tb

def bullets(slide, x, y, w, h, items, size=15, color=WHITE, gap=10):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.space_before = Pt(0)
        r = p.add_run(); r.text = "▸  "
        r.font.size = Pt(size); r.font.color.rgb = GOLD; r.font.bold = True; r.font.name = "Segoe UI"
        r2 = p.add_run(); r2.text = it
        r2.font.size = Pt(size); r2.font.color.rgb = color; r2.font.name = "Segoe UI"
    return tb

def pic_fit(slide, img, x, y, maxw, maxh, border=True):
    iw, ih = Image.open(img).size
    ar = iw / ih
    w = maxw; h = Emu(int(w / ar))
    if h > maxh:
        h = maxh; w = Emu(int(h * ar))
    px = Emu(int(x + (maxw - w) / 2))
    py = Emu(int(y + (maxh - h) / 2))
    if border:
        rect(slide, Emu(px - Emu(0)), py, w, h, fill=None, line=LINE, line_w=1.25)
    slide.shapes.add_picture(img, px, py, w, h)
    return px, py, w, h

# ---------- Slide de capa ----------
s = prs.slides.add_slide(BLANK); bg(s)
rect(s, 0, 0, SW, Inches(0.18), fill=GOLD)
if os.path.exists(ICON):
    s.shapes.add_picture(ICON, Inches(0.9), Inches(2.55), height=Inches(2.0))
text(s, Inches(3.2), Inches(2.35), Inches(9.2), Inches(1.0),
     [("LAB EMC", 22, GOLD, True)], sp_after=2)
text(s, Inches(3.2), Inches(2.95), Inches(9.4), Inches(1.6),
     [("Sistema de Gestão Laboratorial", 40, WHITE, True),
      ("Relatórios CISPR 15  ·  Agenda  ·  Gestão Metrológica", 18, GRAY, False)], sp_after=8)
text(s, Inches(3.2), Inches(5.0), Inches(9.0), Inches(0.6),
     [("LABELO / PUCRS  —  Laboratório de Compatibilidade Eletromagnética", 14, GRAY, False)])
text(s, Inches(0.9), Inches(6.7), Inches(11), Inches(0.5),
     [("Versão 1.0.29  ·  Aplicativo desktop (Windows)  ·  Funciona offline", 12, GOLD, False)])
if os.path.exists(PUCRS):
    s.shapes.add_picture(PUCRS, Inches(11.4), Inches(6.45), height=Inches(0.65))

# ---------- Slide de visão geral ----------
s = prs.slides.add_slide(BLANK); bg(s)
rect(s, 0, 0, SW, Inches(1.15), fill=PANEL)
rect(s, 0, Inches(1.15), SW, Pt(2), fill=GOLD)
text(s, Inches(0.7), Inches(0.28), Inches(12), Inches(0.8),
     [("Visão geral", 30, WHITE, True)])
text(s, Inches(0.7), Inches(1.5), Inches(12), Inches(0.8),
     [("Uma única ferramenta que integra a emissão de relatórios de ensaio CISPR 15 com toda a gestão "
       "metrológica do laboratório — equipamentos, normas, checagens, certificados e procedimentos.", 17, GRAY, False)])
cards = [
    ("Relatórios CISPR 15", "Formulário guiado, geração de PDF e registro automático em Excel."),
    ("Agenda de ensaios", "Controle de protocolos com status de prazo e emissão em lote."),
    ("Gestão metrológica (Lab)", "Equipamentos, normas, checagens, certificados e IT/PC."),
    ("Dados em rede", "Relatórios e emendas reabrem completos em qualquer PC."),
    ("Auto-atualização", "Nova versão chega sozinha, sem reinstalar manualmente."),
    ("Assinatura digital", "PDFs assinados digitalmente com certificado do laboratório."),
]
cw, ch, gx, gy = Inches(3.95), Inches(1.55), Inches(0.25), Inches(0.3)
x0, y0 = Inches(0.7), Inches(2.7)
for i, (t, d) in enumerate(cards):
    cx = x0 + (cw + gx) * (i % 3)
    cy = y0 + (ch + gy) * (i // 3)
    rect(s, cx, cy, cw, ch, fill=PANEL, line=LINE, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, cx + Inches(0.25), cy + Inches(0.2), cw - Inches(0.5), Inches(0.5),
         [(t, 16, GOLD, True)])
    text(s, cx + Inches(0.25), cy + Inches(0.72), cw - Inches(0.5), Inches(0.7),
         [(d, 12.5, GRAY, False)])

# ---------- Slides de módulo ----------
modules = [
    ("Dashboard", "dashboard.png",
     "Visão geral da gestão de EMC", [
        "Cartões com nº de relatórios, emendas, tempos médios e equipamentos cadastrados.",
        "Indicadores de checagens pendentes e vencidas em destaque.",
        "Gráficos: relatórios por mês, não conformidades, checagens por status e equipamentos por grupo.",
        "Ponto de partida diário para a gestão do laboratório.",
     ]),
    ("Relatórios CISPR 15", "cispr15-form.png",
     "Formulário de ensaio de iluminação elétrica", [
        "Seleção de tipo (lâmpada / luminária) e tensões 127 / 220 / 277 V.",
        "Abas Cliente, Emendas e Relatório; leitura de código de barras.",
        "Gera o PDF do relatório e registra automaticamente no Excel.",
        "Emendas e relatórios reabrem completos a partir da rede, em qualquer PC.",
     ]),
    ("Agenda de Execução", "agenda.png",
     "Acompanhamento de protocolos e ensaios", [
        "Status por prazo: Em dia, Vence em ≤ 3 dias e Vencido.",
        "Abas Agenda, Busca, Análise e Follow-up.",
        "Cadastro individual ou em lote, com emissão de lote.",
        "Ordenação por entrada, saída, protocolo ou orçamento.",
     ]),
    ("Equipamentos", "equipamentos.png",
     "Inventário dos instrumentos do laboratório", [
        "Organizados em 6 grupos (Geradores, Medidores, Redes de Impedância, Antenas, Atenuação, Ambientais).",
        "TAG, subgrupo, nº de coleções e status de cada instrumento.",
        "Cadastro, edição e exclusão completos.",
        "Base para checagens, certificados e grandezas metrológicas.",
     ]),
    ("Grupos de Equipamentos", "equipamentos-grupos.png",
     "Taxonomia do acervo metrológico", [
        "Define os 6 grupos e seus subgrupos.",
        "Organiza todo o inventário de instrumentos de forma padronizada.",
        "Facilita filtros, relatórios e templates por subgrupo.",
     ]),
    ("Normas", "normas.png",
     "Biblioteca normativa do laboratório", [
        "CISPR 15, 11 e 32; série IEC 61000-4-x; NBR 15947.",
        "Escopo de cada norma e receptores/medições associados.",
        "Visualizador de PDF local integrado.",
        "Referência direta nos critérios de checagem.",
     ]),
    ("Checagens intermediárias", "checagens-templates.png",
     "Controle de conformidade entre calibrações", [
        "Registro manual, por planilha Excel ou por OCR.",
        "Templates por subgrupo com grandezas e critérios mín./máx. pré-definidos.",
        "Norma de referência vinculada a cada item.",
        "Status por vencimento, alimentando o Dashboard.",
     ]),
    ("Procedimentos (IT / PC)", "procedimentos.png",
     "Checagens e documentação técnica", [
        "Checagem pelo certificado (com interpolação automática) ou manual.",
        "Editor de Instruções de Trabalho e Procedimentos de Calibração no formato LABELO.",
        "Suporte a numeração, tabelas, figuras, definições e referências.",
        "IT/PC casam com a checagem e o certificado.",
     ]),
    ("Grandezas Metrológicas", "grandezas.png",
     "Grandezas por equipamento", [
        "Faixas, unidades e incertezas de cada instrumento.",
        "Alimentam checagens e certificados.",
        "Adição, edição e exclusão por equipamento.",
     ]),
    ("Certificados", "certificados.png",
     "Gestão dos certificados de calibração", [
        "Vincula o certificado ao equipamento e à IT.",
        "Pontos do certificado usados na interpolação da checagem.",
        "Garante rastreabilidade metrológica.",
     ]),
    ("Configurações", "configuracoes.png",
     "Ajustes e integração", [
        "Pasta de dados em rede compartilhada e planilha Excel de registro.",
        "Pasta de cópia de PDF e pasta de atualização (auto-update).",
        "Senha de emissão e certificado digital para assinatura de PDF.",
     ]),
]

for title, img, sub, pts in modules:
    s = prs.slides.add_slide(BLANK); bg(s)
    rect(s, 0, 0, Inches(0.18), SH, fill=GOLD)
    text(s, Inches(0.6), Inches(0.45), Inches(6.2), Inches(0.7),
         [(title, 30, WHITE, True)])
    text(s, Inches(0.62), Inches(1.18), Inches(6.0), Inches(0.6),
         [(sub, 15, GOLD, False)])
    bullets(s, Inches(0.62), Inches(2.1), Inches(5.5), Inches(4.6), pts, size=15.5, gap=14)
    img_path = os.path.join(SHOTS, img)
    if os.path.exists(img_path):
        pic_fit(s, img_path, Inches(6.5), Inches(0.7), Inches(6.5), Inches(6.1))
    # rodapé
    text(s, Inches(0.62), Inches(6.95), Inches(6), Inches(0.4),
         [("CISPR 15 LABELO  ·  v1.0.29", 11, GRAY, False)])

# ---------- Slide final ----------
s = prs.slides.add_slide(BLANK); bg(s)
rect(s, 0, 0, SW, Inches(0.18), fill=GOLD)
text(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.8),
     [("Por que esta ferramenta", 30, WHITE, True)])
diffs = [
    "Aplicativo desktop que roda 100% offline, sem servidor.",
    "Dados em pasta de rede: qualquer PC reabre relatórios e emendas completos.",
    "Auto-atualização — a equipe sempre na versão mais recente, sem reinstalar.",
    "Geração de PDF e registro automático em Excel reduzem retrabalho.",
    "Assinatura digital de PDF com certificado do laboratório.",
    "Importação de checagens por OCR a partir de fotos/planilhas.",
    "Relatórios CISPR 15 e gestão metrológica em uma única interface integrada.",
]
bullets(s, Inches(0.95), Inches(1.9), Inches(11.3), Inches(4.8), diffs, size=18, gap=16)
text(s, Inches(0.95), Inches(6.75), Inches(11), Inches(0.5),
     [("LABELO / PUCRS  ·  Laboratório de Compatibilidade Eletromagnética  ·  v1.0.29", 12, GRAY, False)])

os.makedirs(OUTDIR, exist_ok=True)
prs.save(OUT)
print("SAVED", OUT)
print("SLIDES", len(prs.slides._sldIdLst))
